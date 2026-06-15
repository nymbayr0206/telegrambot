"""
Reusable Python-pptx presentation builder — brand-consistent slides at scale.

Usage:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Define brand colors
    DARK_BG = RGBColor(0x0A, 0x0E, 0x27)
    PRIMARY = RGBColor(0x25, 0x63, 0xEB)
    ACCENT = RGBColor(0x5E, 0xD4, 0xC0)
    GOLD = RGBColor(0xD4, 0xAF, 0x37)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
    CARD = RGBColor(0x14, 0x21, 0x4A)
    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

    # Import helpers from this module
    from presentation_builder import add_bg, add_rect, add_textbox, add_multi_text

    # Build your slides...
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), PRIMARY)
    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
        "Slide Title", 40, WHITE, bold=True)
    prs.save('output.pptx')
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def add_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    """Add a solid-filled rectangle shape (accent bar, card background, etc.)."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text,
                font_size=18, color=None, bold=False,
                alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a single-line or short text box. Returns the text frame for further editing."""
    if color is None:
        color = RGBColor(0xFF, 0xFF, 0xFF)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_multi_text(slide, left, top, width, height, lines,
                   default_size=16, default_color=None, font_name='Calibri',
                   spacing=1.5):
    """Add a textbox with multiple lines, each with optional different formatting.

    Args:
        lines: list of (text, size, color, bold) tuples, or just str for defaults.

    Example:
        add_multi_text(slide, x, y, w, h, [
            ("Main point", 20, WHITE, True),
            ("Sub point", 16, LIGHT, False),
            ("Just a string",),
        ])
    """
    if default_color is None:
        default_color = RGBColor(0xFF, 0xFF, 0xFF)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, size, color, bold = line, default_size, default_color, False
        else:
            text = line[0]
            size = line[1] if len(line) > 1 else default_size
            color = line[2] if len(line) > 2 else default_color
            bold = line[3] if len(line) > 3 else False

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(size * (spacing - 1))
    return tf


def add_card(slide, left, top, width, height,
             title, title_size, title_color,
             body, body_size=14, body_color=None,
             card_color=None, accent_color=None):
    """Add a card (colored rectangle + title + body text)."""
    if card_color is None:
        card_color = RGBColor(0x14, 0x21, 0x4A)
    if body_color is None:
        body_color = RGBColor(0xCC, 0xCC, 0xCC)

    add_rect(slide, left, top, width, height, card_color)
    if accent_color:
        add_rect(slide, left, top, Inches(0.08), height, accent_color)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.2),
                width - Inches(0.6), Inches(0.5),
                title, title_size, title_color, bold=True)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.8),
                width - Inches(0.6), height - Inches(1.0),
                body, body_size, body_color)
